import os
import shutil
import zipfile
import glob
import sqlite3
from core.security import validate_path, audit_log
from pathlib import Path
from langchain_core.tools import Tool, StructuredTool

# New imports for document handling
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("Warning: python-docx not installed. DOCX tools will be unavailable.")

try:
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("Warning: pypdf not installed. PDF tools will be unavailable.")

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False
    print("Warning: openpyxl not installed. Excel tools will be unavailable.")

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    EPUB_AVAILABLE = True
except ImportError:
    EPUB_AVAILABLE = False

try:
    from PIL import Image
    IMAGE_AVAILABLE = True
except ImportError:
    IMAGE_AVAILABLE = False

try:
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

def list_files_tool(directory: str = ".", ws_path: str = ""):
    """Elenca i file nel workspace."""
    target = validate_path(directory, ws_path)
    audit_log("FS_LIST", f"Analisi directory: {directory}")
    return os.listdir(target)

def write_file_tool(filename: str, content: str, ws_path: str = "") -> str:
    """Crea o sovrascrive un file nel workspace in modo sicuro."""
    # Controllo sicurezza contenuto
    dangerous_patterns = ['rm -rf', 'format ', 'eval(', 'exec(', 'os.system', 'subprocess.']
    content_lower = content.lower()
    for pattern in dangerous_patterns:
        if pattern in content_lower:
            audit_log("SECURITY_ALERT", f"Bloccata scrittura file pericoloso: {filename} (pattern: {pattern})")
            return f"Errore: Il contenuto del file è stato bloccato per motivi di sicurezza (rilevato '{pattern}')."

    target = validate_path(filename, ws_path)
    # Assicura che la directory di destinazione esista
    os.makedirs(os.path.dirname(target), exist_ok=True)
    with open(target, "w", encoding="utf-8") as f:
        f.write(content)
    audit_log("FS_WRITE", f"Scrittura file: {filename}")
    return f"Operazione completata nel workspace: {filename} salvato."

def read_file_tool(filename: str, ws_path: str = "") -> str:
    """Legge il contenuto di un file nel workspace."""
    try:
        target = validate_path(filename, ws_path)
        if not os.path.exists(target):
            return f"Errore: Il file {filename} non esiste nel perimetro consentito."
        if os.path.isdir(target):
            return f"Errore: '{filename}' è una directory, non un file. Usa 'list_files' per vederne il contenuto."
        with open(target, "r", encoding="utf-8", errors="replace") as f:
            data = f.read()
        audit_log("FS_READ", f"Lettura file: {filename}")
        return data
    except Exception as e:
        return f"Errore durante la lettura del file: {str(e)}"

def copy_file_tool(src: str, dst: str, ws_path: str = "") -> str:
    """Copia un file o una directory all'interno del workspace."""
    src_path = validate_path(src, ws_path)
    dst_path = validate_path(dst, ws_path)
    os.makedirs(os.path.dirname(dst_path), exist_ok=True)
    if os.path.isdir(src_path):
        shutil.copytree(src_path, dst_path, dirs_exist_ok=True)
    else:
        shutil.copy2(src_path, dst_path)
    audit_log("FS_COPY", f"Copiato: {src} -> {dst}")
    return f"Elemento copiato con successo da {src} a {dst}."

def move_file_tool(src: str, dst: str, ws_path: str = "") -> str:
    """Sposta o rinomina un file o una directory nel workspace."""
    src_path = validate_path(src, ws_path)
    dst_path = validate_path(dst, ws_path)
    os.makedirs(os.path.dirname(dst_path), exist_ok=True)
    shutil.move(src_path, dst_path)
    audit_log("FS_MOVE", f"Spostato: {src} -> {dst}")
    return f"Elemento spostato/rinominato con successo da {src} a {dst}."

def find_files_tool(pattern: str, ws_path: str = "") -> list:
    """Cerca file nel workspace usando un pattern (es. '*.pdf', 'report_*')."""
    target_root = validate_path(".", ws_path)
    search_path = os.path.join(target_root, "**", pattern)
    matches = glob.glob(search_path, recursive=True)
    audit_log("FS_FIND", f"Ricerca pattern: {pattern}")
    return [os.path.relpath(m, target_root) for m in matches]

def delete_file_tool(filename: str, ws_path: str = "") -> str:
    """Elimina permanentemente un file dal workspace."""
    target = validate_path(filename, ws_path)
    audit_log("FS_DELETE", f"AZIONE DISTRUTTIVA: Eliminazione {filename}")
    if os.path.isfile(target):
        os.remove(target)
    elif os.path.isdir(target):
        shutil.rmtree(target)
    return f"L'elemento {filename} è stato rimosso definitivamente."

def read_docx_tool(filename: str, ws_path: str = "") -> str:
    """Legge il contenuto testuale di un file .docx nel workspace."""
    if not DOCX_AVAILABLE:
        return "Errore: La libreria 'python-docx' non è installata. Impossibile leggere file .docx."
    try:
        target = validate_path(filename, ws_path)
        if not os.path.exists(target):
            return f"Errore: Il file {filename} non esiste nel perimetro consentito."
        
        document = Document(target)
        full_text = []
        for para in document.paragraphs:
            full_text.append(para.text)
        audit_log("FS_READ_DOCX", f"Lettura file DOCX: {filename}")
        return "\n".join(full_text)
    except Exception as e:
        return f"Errore durante la lettura del file DOCX: {str(e)}"

def create_docx_tool(filename: str, content: str, ws_path: str = "") -> str:
    """Crea un nuovo file .docx con il contenuto testuale fornito nel workspace."""
    if not DOCX_AVAILABLE:
        return "Errore: La libreria 'python-docx' non è installata. Impossibile creare file .docx."
    try:
        target = validate_path(filename, ws_path)
        os.makedirs(os.path.dirname(target), exist_ok=True)
        
        document = Document()
        for line in content.split('\n'):
            document.add_paragraph(line)
        document.save(target)
        audit_log("FS_CREATE_DOCX", f"Creazione file DOCX: {filename}")
        return f"File DOCX '{filename}' creato con successo."
    except Exception as e:
        return f"Errore durante la creazione del file DOCX: {str(e)}"

def read_pdf_tool(filename: str, ws_path: str = "") -> str:
    """
    Legge il contenuto testuale di un file .pdf nel workspace.
    Se il testo estratto è insufficiente (es. scansione), prova il fallback OCR se disponibile.
    """
    if not PDF_AVAILABLE:
        return "Errore: La libreria 'pypdf' non è installata. Impossibile leggere file .pdf."
    try:
        target = validate_path(filename, ws_path)
        if not os.path.exists(target):
            return f"Errore: Il file {filename} non esiste nel perimetro consentito."
        
        reader = PdfReader(target)
        full_text = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                full_text.append(text)
        
        extracted = "\n".join(full_text).strip()
        
        # Fallback OCR: Se il testo è scarso (< 100 caratteri) e abbiamo i tool necessari
        if len(extracted) < 100 and OCR_AVAILABLE and PDF2IMAGE_AVAILABLE:
            audit_log("FS_PDF_OCR_FALLBACK", f"Avvio OCR su PDF (testo nativo insufficiente): {filename}")
            try:
                # Convertiamo solo le prime 3 pagine per dare un'idea del contenuto senza sovraccaricare il sistema
                images = convert_from_path(target, first_page=1, last_page=3)
                ocr_results = [pytesseract.image_to_string(img) for img in images]
                return "\n".join(ocr_results) + "\n\n(Nota: Testo estratto tramite OCR da scansione immagine)"
            except Exception as ocr_err:
                return f"Il PDF non contiene testo selezionabile e il fallback OCR è fallito: {str(ocr_err)}"

        audit_log("FS_READ_PDF", f"Lettura file PDF (testo nativo): {filename}")
        return extracted if extracted else "Il file PDF sembra vuoto o non contiene testo selezionabile."
    except Exception as e:
        return f"Errore durante la lettura del file PDF: {str(e)}"

def read_xlsx_tool(filename: str, sheet_name: str = None, ws_path: str = "") -> str:
    """Legge il contenuto di un file .xlsx nel workspace e lo restituisce come stringa CSV."""
    if not EXCEL_AVAILABLE:
        return "Errore: La libreria 'openpyxl' non è installata. Impossibile leggere file .xlsx."
    try:
        target = validate_path(filename, ws_path)
        if not os.path.exists(target):
            return f"Errore: Il file {filename} non esiste nel perimetro consentito."
        
        workbook = openpyxl.load_workbook(target)
        if sheet_name:
            if sheet_name not in workbook.sheetnames:
                return f"Errore: Il foglio '{sheet_name}' non trovato nel file '{filename}'. Fogli disponibili: {', '.join(workbook.sheetnames)}"
            sheet = workbook[sheet_name]
        else:
            sheet = workbook.active # Default to active sheet
        
        csv_lines = []
        for row in sheet.iter_rows():
            row_values = [str(cell.value if cell.value is not None else '') for cell in row]
            csv_lines.append(','.join(row_values))
        
        audit_log("FS_READ_XLSX", f"Lettura file XLSX: {filename} (Foglio: {sheet.title})")
        return "\n".join(csv_lines)
    except Exception as e:
        return f"Errore durante la lettura del file XLSX: {str(e)}"

def create_xlsx_tool(filename: str, csv_content: str, sheet_name: str = "Sheet1", ws_path: str = "") -> str:
    """Crea un nuovo file .xlsx da una stringa in formato CSV nel workspace."""
    if not EXCEL_AVAILABLE:
        return "Errore: La libreria 'openpyxl' non è installata. Impossibile creare file .xlsx."
    try:
        target = validate_path(filename, ws_path)
        os.makedirs(os.path.dirname(target), exist_ok=True)
        
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.title = sheet_name
        
        for line in csv_content.strip().split('\n'):
            if line.strip():
                sheet.append(line.split(','))
        
        workbook.save(target)
        audit_log("FS_CREATE_XLSX", f"Creazione file XLSX: {filename} (Foglio: {sheet_name})")
        return f"File XLSX '{filename}' creato con successo."
    except Exception as e:
        return f"Errore durante la creazione del file XLSX: {str(e)}"

def read_pptx_tool(filename: str, ws_path: str = "") -> str:
    """Estrae il testo da una presentazione PowerPoint (.pptx)."""
    if not PPTX_AVAILABLE: return "Errore: 'python-pptx' non installato."
    try:
        target = validate_path(filename, ws_path)
        prs = Presentation(target)
        text_runs = []
        for i, slide in enumerate(prs.slides):
            text_runs.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        audit_log("FS_READ_PPTX", f"Lettura slide: {filename}")
        return "\n".join(text_runs)
    except Exception as e:
        return f"Errore lettura PPTX: {str(e)}"

def create_pptx_tool(filename: str, title: str, bullet_points: list, ws_path: str = "") -> str:
    """Crea una semplice presentazione .pptx con un titolo e una lista di punti."""
    if not PPTX_AVAILABLE: return "Errore: 'python-pptx' non installato."
    try:
        target = validate_path(filename, ws_path)
        prs = Presentation()
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        for p in bullet_points:
            tf.add_paragraph().text = str(p)
        prs.save(target)
        audit_log("FS_CREATE_PPTX", f"Creazione slide: {filename}")
        return f"Presentazione '{filename}' creata."
    except Exception as e:
        return f"Errore creazione PPTX: {str(e)}"

def manage_archive_tool(action: str, archive_name: str, folder_to_zip: str = None, ws_path: str = "") -> str:
    """Gestisce file .zip: 'list' (elenca), 'extract' (estrai), 'create' (comprimi)."""
    try:
        target_archive = validate_path(archive_name, ws_path)
        
        if action == "list":
            with zipfile.ZipFile(target_archive, 'r') as zip_ref:
                return "\n".join(zip_ref.namelist())
        
        elif action == "extract":
            extract_path = os.path.join(ws_path, archive_name.replace(".zip", ""))
            os.makedirs(extract_path, exist_ok=True)
            with zipfile.ZipFile(target_archive, 'r') as zip_ref:
                zip_ref.extractall(extract_path)
            audit_log("FS_ZIP_EXTRACT", f"Estratto: {archive_name}")
            return f"Archivio estratto in: {os.path.basename(extract_path)}"
            
        elif action == "create" and folder_to_zip:
            source_path = validate_path(folder_to_zip, ws_path)
            with zipfile.ZipFile(target_archive, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
                for root, _, files in os.walk(source_path):
                    for file in files:
                        full_path = os.path.join(root, file)
                        rel_path = os.path.relpath(full_path, source_path)
                        zip_ref.write(full_path, rel_path)
            audit_log("FS_ZIP_CREATE", f"Creato: {archive_name}")
            return f"Archivio {archive_name} creato con successo."
            
        return "Azione non valida o parametri mancanti."
    except Exception as e:
        return f"Errore gestione archivio: {str(e)}"

def read_epub_tool(filename: str, ws_path: str = "") -> str:
    """Legge il contenuto testuale di un file .epub."""
    if not EPUB_AVAILABLE: return "Errore: 'ebooklib' o 'beautifulsoup4' non installati."
    try:
        target = validate_path(filename, ws_path)
        book = epub.read_epub(target)
        chapters = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                chapters.append(soup.get_text())
        audit_log("FS_READ_EPUB", f"Lettura libro: {filename}")
        return "\n\n".join(chapters)[:10000] # Limite per non sovraccaricare il contesto
    except Exception as e:
        return f"Errore lettura EPUB: {str(e)}"

def inspect_image_tool(filename: str, ws_path: str = "") -> str:
    """Analizza un'immagine per estrarre metadati (dimensioni, formato, mode)."""
    if not IMAGE_AVAILABLE: return "Errore: 'Pillow' non installato."
    try:
        target = validate_path(filename, ws_path)
        with Image.open(target) as img:
            info = {
                "Formato": img.format,
                "Dimensioni": img.size,
                "Modo": img.mode,
                "Info": img.info
            }
        audit_log("FS_IMAGE_INSPECT", f"Ispezione immagine: {filename}")
        return f"Dettagli immagine {filename}: {str(info)}"
    except Exception as e:
        return f"Errore ispezione immagine: {str(e)}"

def ocr_image_tool(filename: str, ws_path: str = "") -> str:
    """Estrae il testo da un'immagine (.jpg, .png, ecc.) nel workspace usando OCR."""
    if not OCR_AVAILABLE or not IMAGE_AVAILABLE:
        return "Errore: 'pytesseract' o 'Pillow' non installati."
    try:
        target = validate_path(filename, ws_path)
        # Nota: Se tesseract non è nel PATH, aggiungere qui il percorso dell'eseguibile
        text = pytesseract.image_to_string(Image.open(target))
        audit_log("FS_OCR", f"Esecuzione OCR su: {filename}")
        return text if text.strip() else "Nessun testo rilevato nell'immagine."
    except Exception as e:
        return f"Errore OCR: {str(e)}"

def query_sqlite_tool(filename: str, query: str, ws_path: str = "") -> str:
    """Esegue una query SQL (SELECT, INSERT, ecc.) su un database SQLite (.db, .sqlite) nel workspace."""
    try:
        target = validate_path(filename, ws_path)
        if not os.path.exists(target):
            return f"Errore: Il database {filename} non esiste."
        
        conn = sqlite3.connect(target)
        cursor = conn.cursor()
        cursor.execute(query)
        
        if query.strip().upper().startswith("SELECT"):
            rows = cursor.fetchall()
            colnames = [description[0] for description in cursor.description]
            conn.close()
            if not rows:
                return "Query eseguita. Nessun risultato trovato."
            
            # Formattazione base tipo CSV per l'LLM
            res = [", ".join(colnames)]
            for row in rows[:100]: # Limite di sicurezza
                res.append(", ".join(map(str, row)))
            audit_log("FS_DB_QUERY", f"Query SELECT su: {filename}")
            return "\n".join(res)
        else:
            conn.commit()
            conn.close()
            audit_log("FS_DB_EXEC", f"Esecuzione comando SQL su: {filename}")
            return "Comando SQL eseguito con successo."
    except Exception as e:
        return f"Errore SQL: {str(e)}"

def get_fs_tools(workspace_path: str, context_folder: str = None):
    """Genera i tool legandoli al percorso definito nel config."""
    effective_ws = workspace_path
    if context_folder:
        effective_ws = os.path.join(workspace_path, context_folder)
        os.makedirs(effective_ws, exist_ok=True)

    def write_file_wrapper(filename: str, content: str):
        return write_file_tool(filename, content, effective_ws)

    tools = [
        Tool(
            name="list_files", 
            func=lambda d=".": list_files_tool(d, effective_ws), 
            description="Elenca i file. Il root è già il workspace. NON aggiungere 'workspace/' al percorso."
        ),
        StructuredTool.from_function(
            name="write_file",
            func=write_file_wrapper, 
            description="Scrive un file di testo nel workspace. Usa questo solo per file di testo (.txt, .md, .py, .json, .csv, ecc.)."
        ),
        Tool(
            name="read_file", 
            func=lambda f: read_file_tool(f, effective_ws), 
            description="Legge il contenuto di un file di testo nel workspace (anche .yaml, .xml, .html). USA QUESTO TOOL ogni volta che devi verificare se ci sono nuovi messaggi o aggiornamenti in un file usato come chat. Non usarlo per file binari gestiti da altri tool (.docx, .pdf, .xlsx, .pptx, .epub, .zip, .db)."
        ),
        StructuredTool.from_function(
            name="copy_file",
            func=lambda src, dst: copy_file_tool(src, dst, effective_ws),
            description="Copia file o cartelle nel workspace."
        ),
        StructuredTool.from_function(
            name="move_file",
            func=lambda src, dst: move_file_tool(src, dst, effective_ws),
            description="Sposta o rinomina file o cartelle nel workspace."
        ),
        StructuredTool.from_function(
            name="find_files",
            func=lambda pattern: find_files_tool(pattern, effective_ws),
            description="Cerca file nel workspace usando pattern glob (es. '**/*.txt')."
        ),
        Tool(
            name="delete_file", 
            func=lambda f: delete_file_tool(f, effective_ws), 
            description="Elimina file/cartelle (DISTRUTTIVO)."
        )
    ]

    if DOCX_AVAILABLE:
        tools.extend([
            Tool(
                name="read_docx",
                func=lambda f: read_docx_tool(f, effective_ws),
                description="Legge il contenuto testuale di un file Microsoft Word (.docx) nel workspace."
            ),
            StructuredTool.from_function(
                name="create_docx",
                func=lambda filename, content: create_docx_tool(filename, content, effective_ws),
                description="Crea un nuovo file Microsoft Word (.docx) con il testo fornito nel workspace."
            )
        ])
    
    if PDF_AVAILABLE:
        tools.append(
            Tool(
                name="read_pdf",
                func=lambda f: read_pdf_tool(f, effective_ws),
                description="Legge il contenuto testuale di un file PDF (.pdf) nel workspace."
            )
        )
    
    if EXCEL_AVAILABLE:
        tools.extend([
            StructuredTool.from_function(
                name="read_xlsx",
                func=lambda filename, sheet_name=None: read_xlsx_tool(filename, sheet_name, effective_ws),
                description="Legge il contenuto di un file Microsoft Excel (.xlsx) nel workspace e lo restituisce come stringa CSV. Puoi specificare il 'sheet_name'."
            ),
            StructuredTool.from_function(
                name="create_xlsx",
                func=lambda filename, csv_content, sheet_name="Sheet1": create_xlsx_tool(filename, csv_content, sheet_name, effective_ws),
                description="Crea un nuovo file Microsoft Excel (.xlsx) da una stringa in formato CSV nel workspace. Puoi specificare il 'sheet_name'."
            )
        ])
        
    if PPTX_AVAILABLE:
        tools.extend([
            Tool(
                name="read_pptx",
                func=lambda f: read_pptx_tool(f, effective_ws),
                description="Legge il testo da un file PowerPoint (.pptx)."
            ),
            StructuredTool.from_function(
                name="create_pptx",
                func=lambda filename, title, bullet_points: create_pptx_tool(filename, title, bullet_points, effective_ws),
                description="Crea una presentazione PowerPoint (.pptx) con titolo e punti elenco."
            )
        ])

    tools.append(
        StructuredTool.from_function(
            name="manage_archive",
            func=lambda action, archive_name, folder_to_zip=None: manage_archive_tool(action, archive_name, folder_to_zip, effective_ws),
            description="Gestisce file ZIP. Azioni: 'list', 'extract', 'create'. Per 'create' serve 'folder_to_zip'."
        )
    )

    if EPUB_AVAILABLE:
        tools.append(Tool(
            name="read_epub",
            func=lambda f: read_epub_tool(f, effective_ws),
            description="Legge il contenuto testuale di un e-book (.epub)."
        ))

    if IMAGE_AVAILABLE:
        tools.append(Tool(
            name="inspect_image",
            func=lambda f: inspect_image_tool(f, effective_ws),
            description="Analizza un file immagine (.jpg, .png, ecc.) per estrarre metadati tecnici."
        ))

    if OCR_AVAILABLE:
        tools.append(
            Tool(
                name="ocr_image",
                func=lambda f: ocr_image_tool(f, effective_ws),
                description="Estrae testo da un'immagine nel workspace tramite OCR."
            )
        )

    tools.append(
        StructuredTool.from_function(
            name="query_sqlite",
            func=lambda filename, query: query_sqlite_tool(filename, query, effective_ws),
            description="Interroga o modifica un database SQLite (.db, .sqlite) nel workspace tramite query SQL."
        )
    )

    return tools